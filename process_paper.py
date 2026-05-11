import fitz  # pymupdf
import ollama
import os
import json


def extract_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    return " ".join([page.get_text() for page in doc])

    
def extract_text_by_type(path):
    if path.endswith(".pdf"):
        return extract_pdf(path)

    elif path.endswith(".txt"):
        with open(path, "r", encoding="utf-8") as f:
            return f.read()

    elif path.endswith(".docx"):
        from docx import Document
        doc = Document(path)
        return "\n".join([p.text for p in doc.paragraphs])
        
    elif path.endswith(".pptx"):
        from pptx import Presentation

        try:
            prs = Presentation(path)
            text_runs = []

            for i, slide in enumerate(prs.slides, start=1):

                # 👇 SLIDE HEADER
                text_runs.append(f"\n--- SLIDE {i} ---\n")

                # 👇 Slide content
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_runs.append(shape.text.strip())

                # 👇 Speaker notes
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        text_runs.append("\n--- NOTES ---")
                        text_runs.append(notes)

            return "\n".join(text_runs)

        except Exception:
            print(f"Error reading PPTX file: {path}")
            return None
        
    elif path.endswith(".xlsx") or path.endswith(".xls"):
        import pandas as pd

        try:
            df = pd.read_excel(path)
            return df.to_string()
        except Exception as e:
            print(f"Error reading Excel file: {path}")
            return None    

    else:
        return None    


def ask_llama(text):
    prompt = f"""
Return ONLY valid JSON.

STRICT RULES:
- No markdown
- No explanation
- No extra text
- Must be valid JSON
- Do NOT include ```json or ``` blocks

Fields:
- summary (max 5 sentences)
- key_findings (3 bullet points)
- tags (choose from: [ai, biology, climate, physics])

Text:
{text[:6000]}
"""

    response = ollama.chat(
        model='llama3',
        messages=[{'role': 'user', 'content': prompt}]
    )

    return response['message']['content']

PDF_FOLDER = "data"

for root, dirs, files in os.walk(PDF_FOLDER):
    for file in files:

        path = os.path.join(root, file)
        output_path = os.path.join("outputs", file + ".json")

        if os.path.exists(output_path):
            print(f"Skipping (already processed): {file}")
            continue

        print(f"Processing: {file}")

        text = extract_text_by_type(path)

        # ✅ THIS BLOCK MUST BE INSIDE THE LOOP
        if not text or len(text) < 100:
            print(f"Skipping empty or unsupported file: {file}")
            continue

        result = ask_llama(text)

        result = result.replace("```json", "").replace("```", "").strip()

        try:
            json_obj = json.loads(result)
        except:
            print("Skipping bad JSON:", file)
            continue

        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(json_obj, f, indent=2)

        print(f"Saved: {output_path}")